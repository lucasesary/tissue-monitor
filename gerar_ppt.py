"""
gerar_ppt.py — apresentação executiva Tissue Monitor AT1
Design: dark premium — navy profundo + teal + dourado
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paleta ─────────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0D, 0x1B, 0x2A)   # fundo principal
NAVY2    = RGBColor(0x15, 0x2A, 0x40)   # fundo secundário (cards)
TEAL     = RGBColor(0x00, 0xC9, 0xB1)   # destaque principal
GOLD     = RGBColor(0xF5, 0xA6, 0x23)   # destaque secundário
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x8A, 0x9B, 0xB0)   # texto secundário
GRAY_L   = RGBColor(0xD0, 0xDB, 0xE8)   # linha divisória
GREEN    = RGBColor(0x2E, 0xCC, 0x71)   # concluído
ORANGE   = RGBColor(0xE6, 0x7E, 0x22)   # em andamento

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=False):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = fill
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def txt_multi(slide, lines, x, y, w, h, size=16, color=WHITE,
              bold_first=False, line_spacing=0.55):
    """lines: lista de strings; primeira pode ser em bold se bold_first=True"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold_first and i == 0
        run.font.color.rgb = color
        p.space_after  = Pt(int(size * line_spacing * 72 / 72))
    return box


def bg(slide):
    """Fundo escuro padrão."""
    rect(slide, 0, 0, W, H, NAVY)


def top_bar(slide, cor=TEAL):
    """Barra fina no topo."""
    rect(slide, 0, 0, W, Inches(0.07), cor)


def left_bar(slide, cor=TEAL, width=Inches(0.06)):
    rect(slide, 0, 0, width, H, cor)


def rodape(slide, texto="Suzano S.A. — Aracruz  |  AT1 Valmet DCT200HS  |  Confidencial"):
    rect(slide, 0, H - Inches(0.42), W, Inches(0.42), NAVY2)
    txt(slide, texto, Inches(0.4), H - Inches(0.38), W - Inches(0.8), Inches(0.35),
        size=9, color=GRAY_L, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════
def slide_capa(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)

    # bloco diagonal decorativo direito
    rect(slide, W - Inches(4.5), 0, Inches(4.5), H, NAVY2)
    rect(slide, W - Inches(4.6), 0, Inches(0.12), H, TEAL)

    # linha dourada horizontal
    rect(slide, Inches(0.5), Inches(3.6), Inches(5.5), Inches(0.05), GOLD)

    # tag superior
    txt(slide, "SUZANO S.A.  ·  ARACRUZ  ·  2026",
        Inches(0.5), Inches(1.6), Inches(8), Inches(0.5),
        size=11, color=TEAL, bold=True)

    # título principal
    txt(slide, "TISSUE MONITOR",
        Inches(0.5), Inches(2.0), Inches(8.5), Inches(1.2),
        size=52, bold=True, color=WHITE)

    txt(slide, "AT1",
        Inches(0.5), Inches(3.05), Inches(4), Inches(0.9),
        size=52, bold=True, color=TEAL)

    # subtítulo
    txt(slide, "Sistema Inteligente de Monitoramento\ne Suporte à Decisão Operacional",
        Inches(0.5), Inches(3.85), Inches(7.5), Inches(1.1),
        size=18, color=GRAY_L)

    # bloco direito — detalhe técnico
    txt(slide, "MÁQUINA", W - Inches(4.0), Inches(1.8), Inches(3.5), Inches(0.5),
        size=10, color=TEAL, bold=True)
    txt(slide, "Valmet DCT200HS\nViscoNip", W - Inches(4.0), Inches(2.1), Inches(3.5), Inches(0.9),
        size=17, color=WHITE, bold=True)

    txt(slide, "OBJETIVO", W - Inches(4.0), Inches(3.2), Inches(3.5), Inches(0.5),
        size=10, color=TEAL, bold=True)
    txt(slide, "Dados + IA\npara decisões\nmais rápidas", W - Inches(4.0), Inches(3.5),
        Inches(3.5), Inches(1.2), size=17, color=WHITE, bold=True)

    txt(slide, "STATUS", W - Inches(4.0), Inches(5.0), Inches(3.5), Inches(0.5),
        size=10, color=TEAL, bold=True)
    txt(slide, "Fases 0 e 1\nConcluídas ✓", W - Inches(4.0), Inches(5.3), Inches(3.5), Inches(0.9),
        size=17, color=GREEN, bold=True)

    top_bar(slide, TEAL)
    rect(slide, 0, 0, Inches(0.06), H, GOLD)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DIVISOR DE SEÇÃO
# ══════════════════════════════════════════════════════════════════════════
def slide_secao(prs, numero, titulo, subtitulo=""):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    rect(slide, 0, 0, Inches(0.06), H, GOLD)

    # número grande
    txt(slide, f"{numero:02d}", Inches(0.5), Inches(1.5), Inches(3), Inches(2.5),
        size=120, bold=True, color=NAVY2)

    # linha
    rect(slide, Inches(0.5), Inches(3.4), Inches(8), Inches(0.05), TEAL)

    txt(slide, titulo.upper(),
        Inches(0.5), Inches(3.6), Inches(11), Inches(1.2),
        size=42, bold=True, color=WHITE)

    if subtitulo:
        txt(slide, subtitulo, Inches(0.5), Inches(4.8), Inches(9), Inches(0.8),
            size=18, color=GRAY_L, italic=True)

    top_bar(slide)
    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEMA (cards escuros)
# ══════════════════════════════════════════════════════════════════════════
def slide_problema(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, GOLD)

    txt(slide, "SITUAÇÃO ATUAL", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=GOLD, bold=True)
    txt(slide, "5 Gargalos que Custam Produtividade",
        Inches(0.4), Inches(0.5), Inches(11), Inches(0.8),
        size=28, bold=True, color=WHITE)

    rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.04), GRAY_L)

    problemas = [
        ("01", "Dados Isolados",
         "Qualidade, produção e paradas ficam em planilhas Excel no PC local — sem visão integrada."),
        ("02", "Sem Acesso Externo",
         "Operadores não consultam dados fora da planta. Zero mobilidade."),
        ("03", "Análise Manual Lenta",
         "Cruzar processo com qualidade leva horas de trabalho repetitivo."),
        ("04", "Incidentes Repetidos",
         "Sem memória estruturada das causas, os mesmos problemas voltam."),
        ("05", "Conhecimento Disperso",
         "Manuais Valmet e histórico técnico não estão disponíveis na hora H."),
    ]

    x_start = Inches(0.35)
    card_w   = Inches(2.35)
    gap      = Inches(0.18)

    for i, (num, titulo, desc) in enumerate(problemas):
        x = x_start + i * (card_w + gap)
        rect(slide, x, Inches(1.5), card_w, Inches(5.3), NAVY2)
        rect(slide, x, Inches(1.5), card_w, Inches(0.07), GOLD)

        txt(slide, num, x + Inches(0.15), Inches(1.6), card_w, Inches(0.7),
            size=28, bold=True, color=GOLD)
        txt(slide, titulo, x + Inches(0.15), Inches(2.2), card_w - Inches(0.2), Inches(0.7),
            size=14, bold=True, color=WHITE)
        txt(slide, desc, x + Inches(0.15), Inches(2.95), card_w - Inches(0.25), Inches(3.5),
            size=12, color=GRAY_L)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARQUITETURA (fluxo visual)
# ══════════════════════════════════════════════════════════════════════════
def slide_arquitetura(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, TEAL)

    txt(slide, "A SOLUÇÃO", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=TEAL, bold=True)
    txt(slide, "Arquitetura em 4 Camadas",
        Inches(0.4), Inches(0.5), Inches(11), Inches(0.8),
        size=28, bold=True, color=WHITE)

    etapas = [
        (TEAL,   "📧", "E-MAIL",     "Operador exporta e\nenvia a planilha"),
        (GOLD,   "⚙️", "INGESTOR",   "Sistema lê, valida\ne classifica"),
        (RGBColor(0x27, 0x6F, 0xBF), "🗄️", "BANCO",  "Postgres na nuvem\nSão Paulo"),
        (RGBColor(0x8E, 0x44, 0xAD), "🤖", "IA",     "Dashboard +\nAssistente AI"),
    ]

    bw = Inches(2.8)
    bh = Inches(4.0)
    y0 = Inches(1.5)
    gap = Inches(0.38)
    x0 = Inches(0.45)

    for i, (cor, icone, label, desc) in enumerate(etapas):
        x = x0 + i * (bw + gap)
        rect(slide, x, y0, bw, bh, NAVY2)
        rect(slide, x, y0, bw, Inches(0.07), cor)

        txt(slide, icone, x + Inches(0.2), y0 + Inches(0.2), bw, Inches(0.7),
            size=28, color=WHITE)
        txt(slide, label, x + Inches(0.2), y0 + Inches(0.9), bw - Inches(0.3),
            Inches(0.55), size=15, bold=True, color=cor)
        txt(slide, desc, x + Inches(0.2), y0 + Inches(1.5), bw - Inches(0.3),
            Inches(2.2), size=13, color=GRAY_L)

        # seta entre blocos
        if i < len(etapas) - 1:
            ax = x + bw + Inches(0.05)
            txt(slide, "→", ax, y0 + Inches(1.7), Inches(0.3), Inches(0.6),
                size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # benefício abaixo
    rect(slide, Inches(0.45), Inches(5.7), Inches(12.4), Inches(1.1), NAVY2)
    rect(slide, Inches(0.45), Inches(5.7), Inches(0.06), Inches(1.1), TEAL)
    txt(slide, "Resultado: operadores enviam 1 e-mail → dados disponíveis em todo lugar, em minutos, sem depender da rede Suzano.",
        Inches(0.65), Inches(5.82), Inches(12.0), Inches(0.85),
        size=14, color=WHITE, italic=True)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DADOS NA NUVEM (KPIs grandes)
# ══════════════════════════════════════════════════════════════════════════
def slide_dados(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, TEAL)

    txt(slide, "BANCO DE DADOS — NEON / SÃO PAULO", Inches(0.4), Inches(0.15),
        Inches(12), Inches(0.5), size=11, color=TEAL, bold=True)
    txt(slide, "Dados de Maio/2026 — Já na Nuvem",
        Inches(0.4), Inches(0.5), Inches(11), Inches(0.8),
        size=28, bold=True, color=WHITE)

    kpis = [
        ("1.619", "jumbos de\nqualidade", TEAL,
         "Gramatura · Espessura\nHandfeel · Tração"),
        ("1.626", "jumbos de\nprodução", GOLD,
         "Velocidade · Quebras\nFamília · Duração"),
        ("2.450", "eventos de\ndowntime", RGBColor(0xE7, 0x4C, 0x3C),
         "Classe · Tipo · Causa\nDuração em minutos"),
        ("34", "manuais\nValmet", RGBColor(0x8E, 0x44, 0xAD),
         "Indexados para\no assistente AI"),
    ]

    kw = Inches(2.9)
    kh = Inches(4.5)
    y0 = Inches(1.5)
    gap = Inches(0.27)
    x0 = Inches(0.45)

    for i, (num, label, cor, detalhe) in enumerate(kpis):
        x = x0 + i * (kw + gap)
        rect(slide, x, y0, kw, kh, NAVY2)
        rect(slide, x, y0 + kh - Inches(0.07), kw, Inches(0.07), cor)

        txt(slide, num, x + Inches(0.15), y0 + Inches(0.3), kw - Inches(0.2),
            Inches(1.3), size=52, bold=True, color=cor, align=PP_ALIGN.CENTER)
        txt(slide, label, x + Inches(0.15), y0 + Inches(1.55), kw - Inches(0.2),
            Inches(0.8), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, x + Inches(0.4), y0 + Inches(2.45), kw - Inches(0.8), Inches(0.03), GRAY_L)
        txt(slide, detalhe, x + Inches(0.15), y0 + Inches(2.6), kw - Inches(0.2),
            Inches(1.5), size=12, color=GRAY_L, align=PP_ALIGN.CENTER)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD
# ══════════════════════════════════════════════════════════════════════════
def slide_dashboard(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, GOLD)

    txt(slide, "DASHBOARD", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=GOLD, bold=True)
    txt(slide, "Acesso em Qualquer Lugar — Celular, Tablet ou PC",
        Inches(0.4), Inches(0.5), Inches(12), Inches(0.8),
        size=28, bold=True, color=WHITE)

    # coluna esquerda — abas
    abas = [
        (TEAL,   "Processo",   "Curvas OPC em tempo quase real"),
        (GOLD,   "Qualidade",  "Carta de Controle IMR + conformidade por produto"),
        (RGBColor(0xE7, 0x4C, 0x3C), "Downtime", "Pareto de causas + Top 3 paradas + MCR"),
        (RGBColor(0x27, 0x6F, 0xBF), "Proc × Qual", "Mapa de correlações + influenciadores"),
        (RGBColor(0x8E, 0x44, 0xAD), "Yankee",    "Temperaturas LA / Meio / LC"),
        (GREEN,  "Relatórios", "PDF semana × semana + análises salvas"),
    ]

    y = Inches(1.5)
    for cor, nome, desc in abas:
        rect(slide, Inches(0.4), y + Inches(0.08), Inches(0.06), Inches(0.52), cor)
        rect(slide, Inches(0.5), y, Inches(6.8), Inches(0.68), NAVY2)
        txt(slide, nome, Inches(0.65), y + Inches(0.04), Inches(2.2), Inches(0.55),
            size=13, bold=True, color=cor)
        txt(slide, desc, Inches(2.6), y + Inches(0.12), Inches(4.6), Inches(0.5),
            size=12, color=GRAY_L)
        y += Inches(0.76)

    # coluna direita — acesso
    rect(slide, Inches(7.7), Inches(1.5), Inches(5.2), Inches(5.3), NAVY2)
    rect(slide, Inches(7.7), Inches(1.5), Inches(5.2), Inches(0.07), TEAL)

    txt(slide, "ACESSO E SEGURANÇA", Inches(7.9), Inches(1.65), Inches(4.8), Inches(0.5),
        size=11, color=TEAL, bold=True)

    items = [
        "🌐  Link HTTPS — sem instalação",
        "📱  Celular, tablet ou PC de casa",
        "🔐  Login por e-mail, sem VPN",
        "🔄  Atualiza quando chega novo dado",
        "🏭  Pronto para múltiplas fábricas",
    ]
    y2 = Inches(2.25)
    for item in items:
        txt(slide, item, Inches(7.9), y2, Inches(4.8), Inches(0.55), size=14, color=WHITE)
        y2 += Inches(0.7)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ASSISTENTE AI
# ══════════════════════════════════════════════════════════════════════════
def slide_ai(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, TEAL)

    txt(slide, "ASSISTENTE AI", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=TEAL, bold=True)
    txt(slide, "Especialista Disponível 24h para os Operadores",
        Inches(0.4), Inches(0.5), Inches(12), Inches(0.8),
        size=28, bold=True, color=WHITE)

    # simulação de conversa — lado esquerdo
    rect(slide, Inches(0.4), Inches(1.5), Inches(6.5), Inches(5.3), NAVY2)
    rect(slide, Inches(0.4), Inches(1.5), Inches(6.5), Inches(0.07), TEAL)

    txt(slide, "EXEMPLO DE CONVERSA", Inches(0.6), Inches(1.65), Inches(6.0), Inches(0.4),
        size=10, color=TEAL, bold=True)

    # balão operador
    rect(slide, Inches(0.6), Inches(2.15), Inches(5.8), Inches(0.9), RGBColor(0x1E, 0x3A, 0x52))
    txt(slide, "👷  Operador:  \"Espessura caindo desde 22h.\n       O que devo verificar primeiro?\"",
        Inches(0.75), Inches(2.2), Inches(5.5), Inches(0.8), size=12, color=WHITE)

    # balão IA
    rect(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(2.4), RGBColor(0x0A, 0x3D, 0x3A))
    rect(slide, Inches(0.6), Inches(3.2), Inches(0.05), Inches(2.4), TEAL)
    txt(slide, "🤖  Assistente:  Analisei os dados das últimas 4h.\n"
               "       Pressão do ViscoNip caiu 3,2 bar (22h05).\n"
               "       Manual Valmet §4.3: verifique válvula V12.\n\n"
               "       Quer que eu abra o histórico de V12?",
        Inches(0.75), Inches(3.3), Inches(5.5), Inches(2.2), size=12, color=WHITE)

    # lado direito — como funciona
    rect(slide, Inches(7.3), Inches(1.5), Inches(5.6), Inches(5.3), NAVY2)
    rect(slide, Inches(7.3), Inches(1.5), Inches(5.6), Inches(0.07), GOLD)

    txt(slide, "COMO FUNCIONA", Inches(7.5), Inches(1.65), Inches(5.2), Inches(0.4),
        size=10, color=GOLD, bold=True)

    passos = [
        ("1", "Operador descreve o problema em linguagem natural"),
        ("2", "IA cruza dados OPC + qualidade + downtime do período"),
        ("3", "Busca nos 34 manuais Valmet a recomendação técnica"),
        ("4", "Faz perguntas e guia o diagnóstico passo a passo"),
        ("5", "Cada conversa vira aprendizado — sistema evolui"),
    ]

    y3 = Inches(2.15)
    for num, passo in passos:
        rect(slide, Inches(7.5), y3 + Inches(0.07), Inches(0.38), Inches(0.38),
             GOLD)
        txt(slide, num, Inches(7.5), y3, Inches(0.45), Inches(0.55),
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(slide, passo, Inches(8.05), y3 + Inches(0.04), Inches(4.6), Inches(0.5),
            size=13, color=WHITE)
        y3 += Inches(0.85)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════
def slide_roadmap(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, GOLD)

    txt(slide, "ROADMAP", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=GOLD, bold=True)
    txt(slide, "Plano de Implementação — 5 Fases",
        Inches(0.4), Inches(0.5), Inches(11), Inches(0.8),
        size=28, bold=True, color=WHITE)

    # linha do tempo
    rect(slide, Inches(0.7), Inches(3.2), Inches(12.0), Inches(0.06), GRAY_L)

    fases = [
        (GREEN,   "✓", "FASE 0",   "Parsers\nRefatorados",   "Concluída"),
        (GREEN,   "✓", "FASE 1",   "Banco\nPostgres",        "Concluída"),
        (ORANGE,  "►", "FASE 2",   "Ingestor\nde E-mail",    "Em andamento"),
        (NAVY2,   "○", "FASE 3",   "Dashboard\nExterno",     "Próxima"),
        (NAVY2,   "○", "FASE 4",   "Assistente\nAI",         "Planejada"),
    ]

    fw = Inches(2.35)
    gap = Inches(0.13)
    x0 = Inches(0.7)

    for i, (cor, icone, fase, desc, status) in enumerate(fases):
        x = x0 + i * (fw + gap)

        # círculo na linha
        cx = x + fw / 2 - Inches(0.3)
        rect(slide, cx, Inches(3.0), Inches(0.6), Inches(0.6), cor)
        txt(slide, icone, cx, Inches(2.98), Inches(0.6), Inches(0.62),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # card
        rect(slide, x, Inches(1.45), fw, Inches(1.45), NAVY2)
        rect(slide, x, Inches(1.45), fw, Inches(0.06), cor)

        txt(slide, fase, x + Inches(0.15), Inches(1.55), fw - Inches(0.2), Inches(0.45),
            size=10, bold=True, color=cor)
        txt(slide, desc, x + Inches(0.15), Inches(1.95), fw - Inches(0.2), Inches(0.8),
            size=13, bold=True, color=WHITE)

        # status abaixo
        rect(slide, x, Inches(3.75), fw, Inches(1.4), NAVY2)
        txt(slide, status, x + Inches(0.15), Inches(3.85), fw - Inches(0.2), Inches(0.6),
            size=12, color=cor, bold=True)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ESCALABILIDADE
# ══════════════════════════════════════════════════════════════════════════
def slide_escalabilidade(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    top_bar(slide)
    left_bar(slide, TEAL)

    txt(slide, "MODELO DE NEGÓCIO", Inches(0.4), Inches(0.15), Inches(10), Inches(0.5),
        size=11, color=TEAL, bold=True)
    txt(slide, "Validado na AT1 — Vendável para Qualquer Fábrica",
        Inches(0.4), Inches(0.5), Inches(12), Inches(0.8),
        size=28, bold=True, color=WHITE)

    diferenciais = [
        (TEAL,   "Zero Instalação",
                 "Funciona 100% via e-mail e navegador.\nNenhum software na rede do cliente."),
        (GOLD,   "Multi-Fábrica",
                 "Cada cliente tem seus dados isolados.\nEscala para N fábricas sem reescrever."),
        (GREEN,  "IA Embarcada",
                 "Manuais técnicos + histórico da máquina\nintegrados ao assistente."),
        (RGBColor(0x8E, 0x44, 0xAD), "Baixo Custo",
                 "Infraestrutura cloud sob demanda.\nCusto cresce só com uso."),
    ]

    dw = Inches(2.95)
    dh = Inches(4.2)
    gap = Inches(0.2)
    x0 = Inches(0.4)
    y0 = Inches(1.5)

    for i, (cor, titulo, desc) in enumerate(diferenciais):
        x = x0 + i * (dw + gap)
        rect(slide, x, y0, dw, dh, NAVY2)
        rect(slide, x, y0, dw, Inches(0.07), cor)

        txt(slide, titulo, x + Inches(0.2), y0 + Inches(0.25), dw - Inches(0.3),
            Inches(0.65), size=15, bold=True, color=cor)
        rect(slide, x + Inches(0.2), y0 + Inches(0.9), dw - Inches(0.4), Inches(0.03), GRAY_L)
        txt(slide, desc, x + Inches(0.2), y0 + Inches(1.05), dw - Inches(0.3),
            Inches(2.8), size=13, color=GRAY_L)

    rodape(slide)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ENCERRAMENTO / CTA
# ══════════════════════════════════════════════════════════════════════════
def slide_cta(prs):
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    rect(slide, 0, 0, Inches(0.06), H, TEAL)

    rect(slide, W - Inches(5.0), 0, Inches(5.0), H, NAVY2)
    rect(slide, W - Inches(5.05), 0, Inches(0.07), H, TEAL)

    # esquerda
    txt(slide, "PRÓXIMO PASSO", Inches(0.5), Inches(1.5), Inches(7), Inches(0.6),
        size=12, color=TEAL, bold=True)
    txt(slide, "Fase 2 —\nIngestor de E-mail",
        Inches(0.5), Inches(2.0), Inches(7.5), Inches(2.0),
        size=44, bold=True, color=WHITE)
    rect(slide, Inches(0.5), Inches(4.0), Inches(5.5), Inches(0.06), GOLD)
    txt(slide, "Meta: primeiros dados automáticos\nfluindo ainda esta semana.",
        Inches(0.5), Inches(4.2), Inches(7.0), Inches(1.2),
        size=18, color=GRAY_L)

    # direita
    txt(slide, "CONTATO", W - Inches(4.6), Inches(2.0), Inches(4.2), Inches(0.5),
        size=10, color=TEAL, bold=True)
    txt(slide, "Lucas Brígido", W - Inches(4.6), Inches(2.4), Inches(4.2), Inches(0.7),
        size=22, bold=True, color=WHITE)
    txt(slide, "Suzano S.A. — Aracruz", W - Inches(4.6), Inches(3.0), Inches(4.2), Inches(0.5),
        size=14, color=GRAY_L)
    txt(slide, "lucasesary@gmail.com", W - Inches(4.6), Inches(3.5), Inches(4.2), Inches(0.5),
        size=13, color=TEAL)

    txt(slide, "AT1 Valmet DCT200HS", W - Inches(4.6), Inches(4.5), Inches(4.2), Inches(0.5),
        size=11, color=GRAY_L, italic=True)

    top_bar(slide)


# ══════════════════════════════════════════════════════════════════════════
# GERAR
# ══════════════════════════════════════════════════════════════════════════
slide_capa(prs)
slide_secao(prs, 1, "O Problema", "Por que precisamos mudar a forma de trabalhar")
slide_problema(prs)
slide_secao(prs, 2, "A Solução", "Arquitetura completa do Tissue Monitor")
slide_arquitetura(prs)
slide_secao(prs, 3, "Dados na Nuvem", "O que já está pronto e funcionando")
slide_dados(prs)
slide_secao(prs, 4, "Dashboard & AI", "O que os operadores vão acessar")
slide_dashboard(prs)
slide_ai(prs)
slide_secao(prs, 5, "Roadmap", "Da AT1 para o mercado")
slide_roadmap(prs)
slide_escalabilidade(prs)
slide_cta(prs)

SAIDA = "Tissue_Monitor_AT1_v2.pptx"
prs.save(SAIDA)
print(f"PPT salvo: {SAIDA}")
print(f"Total de slides: {len(prs.slides)}")
